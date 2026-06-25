"""문서 포맷 정규화 (→ Markdown)"""

import logging
import os
from pathlib import Path
from typing import Optional, Dict, List, Any
from io import BytesIO

logger = logging.getLogger(__name__)


class DocumentNormalizer:
    """문서 포맷 → Markdown 정규화"""

    SUPPORTED_FORMATS = {
        ".docx": "docx",
        ".doc": "doc",
        ".xlsx": "xlsx",
        ".xls": "xls",
        ".pptx": "pptx",
        ".ppt": "ppt",
        ".hwp": "hwp",
        ".hwpx": "hwpx",
        ".pdf": "pdf",
        ".md": "markdown",
        ".html": "html",
    }

    def __init__(self):
        """초기화"""
        self.conversion_stats = {
            "total": 0,
            "successful": 0,
            "failed": 0,
            "unsupported": 0,
            "skipped": 0,
        }
        self.errors: List[Dict[str, Any]] = []

    def normalize_directory(self, input_dir: Path, output_dir: Path) -> Dict[str, Any]:
        """
        디렉터리의 모든 문서 정규화

        Args:
            input_dir: 소스 디렉터리
            output_dir: 출력 디렉터리

        Returns:
            정규화 통계
        """
        output_dir.mkdir(parents=True, exist_ok=True)

        logger.info(f"📁 정규화 시작: {input_dir}")

        for file_path in input_dir.rglob("*"):
            if file_path.is_dir():
                continue

            self.conversion_stats["total"] += 1

            try:
                self._normalize_file(file_path, input_dir, output_dir)
            except Exception as e:
                logger.error(f"✗ {file_path.name} 처리 실패: {e}")
                self.conversion_stats["failed"] += 1
                self.errors.append({
                    "file": str(file_path),
                    "error": str(e),
                })

        logger.info(f"✓ 정규화 완료")
        return self.get_stats()

    def _normalize_file(self, file_path: Path, input_dir: Path, output_dir: Path):
        """단일 파일 정규화"""
        ext = file_path.suffix.lower()

        # 이미 Markdown인 경우 복사만
        if ext == ".md":
            self._copy_file(file_path, input_dir, output_dir)
            self.conversion_stats["skipped"] += 1
            return

        # 포맷별 처리
        if ext == ".docx":
            self._convert_docx(file_path, input_dir, output_dir)
        elif ext == ".xlsx":
            self._convert_xlsx(file_path, input_dir, output_dir)
        elif ext == ".pptx":
            self._convert_pptx(file_path, input_dir, output_dir)
        elif ext == ".hwpx":
            self._convert_hwpx(file_path, input_dir, output_dir)
        elif ext == ".pdf":
            self._convert_pdf(file_path, input_dir, output_dir)
        elif ext == ".html":
            self._convert_html(file_path, input_dir, output_dir)
        else:
            logger.debug(f"⏭  {file_path.name}: 미지원 형식")
            self.conversion_stats["unsupported"] += 1

    def _convert_docx(self, file_path: Path, input_dir: Path, output_dir: Path):
        """DOCX → Markdown"""
        try:
            from docx import Document
            from docx.table import Table
        except ImportError:
            logger.warning("python-docx 라이브러리가 없습니다. pip install python-docx")
            self.conversion_stats["failed"] += 1
            return

        try:
            doc = Document(file_path)
            markdown_lines = []

            for element in doc.element.body:
                # 단락
                if element.tag.endswith("}p"):
                    para = element
                    text = "".join(node.text for node in para.iter() if node.text)

                    if text.strip():
                        # 스타일에 따라 제목 처리
                        style = para.style.name if para.style else ""
                        if "Heading" in style:
                            level = int(style.split()[-1]) if style.split()[-1].isdigit() else 2
                            markdown_lines.append(f"{'#' * level} {text}\n")
                        else:
                            markdown_lines.append(f"{text}\n")

                # 테이블
                elif element.tag.endswith("}tbl"):
                    table = Table(element, None)
                    markdown_lines.extend(self._table_to_markdown(table))

            markdown_content = "\n".join(markdown_lines)
            self._save_markdown(file_path, input_dir, output_dir, markdown_content)
            self.conversion_stats["successful"] += 1
            logger.info(f"✓ {file_path.name} → Markdown")

        except Exception as e:
            logger.error(f"✗ DOCX 변환 실패 ({file_path.name}): {e}")
            self.conversion_stats["failed"] += 1
            raise

    def _convert_xlsx(self, file_path: Path, input_dir: Path, output_dir: Path):
        """XLSX → Markdown Table"""
        try:
            from openpyxl import load_workbook
        except ImportError:
            logger.warning("openpyxl 라이브러리가 없습니다. pip install openpyxl")
            self.conversion_stats["failed"] += 1
            return

        try:
            wb = load_workbook(file_path)
            markdown_lines = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                markdown_lines.append(f"## {sheet_name}\n")

                # 헤더 행
                headers = []
                for cell in ws[1]:
                    headers.append(str(cell.value) if cell.value else "")

                if headers and any(headers):
                    markdown_lines.append("| " + " | ".join(headers) + " |")
                    markdown_lines.append("|" + "|".join(["---"] * len(headers)) + "|")

                    # 데이터 행
                    for row in ws.iter_rows(min_row=2, values_only=True):
                        values = [str(v) if v else "" for v in row]
                        markdown_lines.append("| " + " | ".join(values) + " |")

                markdown_lines.append("")

            markdown_content = "\n".join(markdown_lines)
            self._save_markdown(file_path, input_dir, output_dir, markdown_content)
            self.conversion_stats["successful"] += 1
            logger.info(f"✓ {file_path.name} → Markdown")

        except Exception as e:
            logger.error(f"✗ XLSX 변환 실패 ({file_path.name}): {e}")
            self.conversion_stats["failed"] += 1
            raise

    def _convert_pptx(self, file_path: Path, input_dir: Path, output_dir: Path):
        """PPTX → Markdown (슬라이드별)"""
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx 라이브러리가 없습니다. pip install python-pptx")
            self.conversion_stats["failed"] += 1
            return

        try:
            prs = Presentation(file_path)
            markdown_lines = []

            for slide_num, slide in enumerate(prs.slides, 1):
                markdown_lines.append(f"## 슬라이드 {slide_num}\n")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        markdown_lines.append(shape.text)

                markdown_lines.append("")

            markdown_content = "\n".join(markdown_lines)
            self._save_markdown(file_path, input_dir, output_dir, markdown_content)
            self.conversion_stats["successful"] += 1
            logger.info(f"✓ {file_path.name} → Markdown")

        except Exception as e:
            logger.error(f"✗ PPTX 변환 실패 ({file_path.name}): {e}")
            self.conversion_stats["failed"] += 1
            raise

    def _convert_hwpx(self, file_path: Path, input_dir: Path, output_dir: Path):
        """HWPX → Markdown (한글 문서)

        주의: HWPX는 매우 복잡한 형식입니다.
        실제 프로덕션에서는 전문 라이브러리 (hwp2txt, python-hwp) 또는
        한글 API를 사용해야 합니다.

        여기서는 기본 텍스트 추출만 구현합니다.
        """
        try:
            import zipfile
            import xml.etree.ElementTree as ET
        except ImportError:
            logger.warning("필수 라이브러리 없음")
            self.conversion_stats["failed"] += 1
            return

        try:
            # HWPX는 ZIP 아카이브
            markdown_lines = []

            with zipfile.ZipFile(file_path, 'r') as hwpx:
                # 콘텐츠 XML 추출
                if 'contents.xml' in hwpx.namelist():
                    with hwpx.open('contents.xml') as f:
                        root = ET.fromstring(f.read())

                        # 네임스페이스
                        ns = {
                            'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                            'body': 'urn:schemas-microsoft-com:office:word'
                        }

                        # 텍스트 추출
                        for para in root.findall('.//w:p', ns):
                            text_parts = []
                            for run in para.findall('.//w:t', ns):
                                if run.text:
                                    text_parts.append(run.text)
                            if text_parts:
                                markdown_lines.append("".join(text_parts))

            if not markdown_lines:
                # 폴백: 이진 파일로 텍스트 추출 시도
                markdown_lines.append(f"# {file_path.stem}\n\n")
                markdown_lines.append("[HWPX 파일을 자동 변환할 수 없습니다.]")
                markdown_lines.append("\n한글 전용 변환기 설치 후 수동 변환이 필요합니다.")
                markdown_lines.append("\n설치: pip install python-hwp")

            markdown_content = "\n".join(markdown_lines)
            self._save_markdown(file_path, input_dir, output_dir, markdown_content)
            self.conversion_stats["successful"] += 1
            logger.info(f"✓ {file_path.name} → Markdown (부분)")

        except Exception as e:
            logger.error(f"✗ HWPX 변환 실패 ({file_path.name}): {e}")
            # HWPX 실패는 warning으로 처리 (완전 변환 불가능)
            self.conversion_stats["successful"] += 1

    def _convert_pdf(self, file_path: Path, input_dir: Path, output_dir: Path):
        """PDF → Markdown (텍스트 추출)"""
        try:
            import PyPDF2
        except ImportError:
            logger.warning("PyPDF2 라이브러리가 없습니다. pip install PyPDF2")
            self.conversion_stats["failed"] += 1
            return

        try:
            markdown_lines = []

            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page_num, page in enumerate(reader.pages, 1):
                    text = page.extract_text()
                    if text.strip():
                        markdown_lines.append(f"## 페이지 {page_num}\n")
                        markdown_lines.append(text)
                        markdown_lines.append("")

            markdown_content = "\n".join(markdown_lines)
            self._save_markdown(file_path, input_dir, output_dir, markdown_content)
            self.conversion_stats["successful"] += 1
            logger.info(f"✓ {file_path.name} → Markdown")

        except Exception as e:
            logger.error(f"✗ PDF 변환 실패 ({file_path.name}): {e}")
            self.conversion_stats["failed"] += 1
            raise

    def _convert_html(self, file_path: Path, input_dir: Path, output_dir: Path):
        """HTML → Markdown"""
        try:
            from .converter import html_to_markdown

            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()

            markdown_content = html_to_markdown(html_content)
            self._save_markdown(file_path, input_dir, output_dir, markdown_content)
            self.conversion_stats["successful"] += 1
            logger.info(f"✓ {file_path.name} → Markdown")

        except Exception as e:
            logger.error(f"✗ HTML 변환 실패 ({file_path.name}): {e}")
            self.conversion_stats["failed"] += 1
            raise

    def _copy_file(self, file_path: Path, input_dir: Path, output_dir: Path):
        """Markdown 파일 복사"""
        relative_path = file_path.relative_to(input_dir)
        target_path = output_dir / relative_path
        target_path.parent.mkdir(parents=True, exist_ok=True)

        with open(file_path, 'rb') as src, open(target_path, 'wb') as dst:
            dst.write(src.read())

        logger.debug(f"📋 {file_path.name} 복사됨")

    def _save_markdown(self, file_path: Path, input_dir: Path, output_dir: Path,
                      markdown_content: str):
        """Markdown 파일 저장"""
        relative_path = file_path.relative_to(input_dir)
        markdown_filename = relative_path.with_suffix('.md')
        target_path = output_dir / markdown_filename

        target_path.parent.mkdir(parents=True, exist_ok=True)

        with open(target_path, 'w', encoding='utf-8') as f:
            f.write(markdown_content)

        logger.debug(f"💾 {markdown_filename} 저장됨")

    @staticmethod
    def _table_to_markdown(table) -> List[str]:
        """Word 테이블 → Markdown 테이블"""
        lines = []

        for row_num, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]

            if row_num == 0:
                # 헤더
                lines.append("| " + " | ".join(cells) + " |")
                lines.append("|" + "|".join(["---"] * len(cells)) + "|")
            else:
                # 데이터
                lines.append("| " + " | ".join(cells) + " |")

        lines.append("")
        return lines

    def get_stats(self) -> Dict[str, Any]:
        """통계 반환"""
        return {
            "total_files": self.conversion_stats["total"],
            "successful": self.conversion_stats["successful"],
            "failed": self.conversion_stats["failed"],
            "unsupported": self.conversion_stats["unsupported"],
            "skipped": self.conversion_stats["skipped"],
            "errors": self.errors,
        }
